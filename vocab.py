#!/usr/bin/env python
# -*- coding: utf-8 -*-

from pptx import Presentation
import pandas as pd
import numpy as np

def getVocab(path_name):
    prs = Presentation(path_name)
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []

    for slide in prs.slides:
        slide_info = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            slide_info.append(shape.text)
        text_runs.append(slide_info)
    return text_runs[1:]

def filterVocab(vocab):
    # remove the weird thing
    vocab = [[item for item in slide if "I-" not in item] for slide in vocab]

    # remove pinyin from list
    tones = ['ā', 'ē', 'ī', 'ō', 'ū', 'ǖ', 'á', 'é', 'í', 'ó', 'ú', 'ǘ', 'ǎ', 'ě', 'ǐ', 'ǒ', 'ǔ', 'ǚ', 'à', 'è', 'ì', 'ò', 'ù', 'ǜ']
    pinyin = [text for slide in vocab for text in slide if set(text) & set(tones)]
    vocab = [[item for item in slide if item not in pinyin] for slide in vocab]

    # get description (assuming longer string length is pinyin)
    description = [max(slide, key=len) for slide in vocab]
    hanzi = [item for slide in vocab for item in slide if item not in description]

    return hanzi, pinyin, description

def main():
    mhanzi, mpinyin, mdescription = [], [], []

    # loop through all files and put data into lists
    for lesson in range(1, 10):
        path_name = r"C:\Users\Kieran\Documents\Projects\mandarin\DangDai B6  e.g. of Grammar & vocabulary  (traditional)\Dangdai_6_Vocabulary\B6-L0{}.pptx".format(str(lesson))
        vocab = getVocab(path_name)
        hanzi, pinyin, description = filterVocab(vocab)
        mhanzi += hanzi; mpinyin += pinyin; mdescription += description

    vocab = getVocab(r"C:\Users\Kieran\Documents\Projects\mandarin\DangDai B6  e.g. of Grammar & vocabulary  (traditional)\Dangdai_6_Vocabulary\B6-L10.pptx")
    hanzi, pinyin, description = filterVocab(vocab)
    mhanzi += hanzi; mpinyin += pinyin; mdescription += description

    # put everything in the df, needs to be transposed for anki
    df = pd.DataFrame(data=np.array([hanzi, pinyin, description]).transpose())
    df.head()

    # export to csv (which is actually a tsv but oh well)
    df.to_csv('file_name.csv', encoding='utf-8', index=False, header=False, sep='\t')
    return

if __name__ == "__main__":
    main()
